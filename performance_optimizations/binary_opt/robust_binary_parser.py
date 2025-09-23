"""
Robust Binary Format Parser
Enhanced VISIO/PowerPoint parsing with error handling and fallback mechanisms.
"""

import io
import json
import logging
import mimetypes
import os
import tempfile
import traceback
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, Union
import zipfile
import xml.etree.ElementTree as ET

# Third-party imports with fallback handling
try:
    from pptx import Presentation
    from pptx.exc import PackageNotFoundError, PresentationError
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx not available. PowerPoint parsing will be limited.")

try:
    import olefile
    OLE_AVAILABLE = True
except ImportError:
    OLE_AVAILABLE = False
    logging.warning("olefile not available. Legacy Office format support limited.")

try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False
    logging.warning("PIL not available. Image extraction will be limited.")


class BinaryFormatError(Exception):
    """Custom exception for binary format parsing errors."""
    pass


class FormatValidator:
    """Validates and identifies binary file formats."""
    
    # File signatures for format detection
    SIGNATURES = {
        # Office formats
        b'\x50\x4B\x03\x04': 'office_xml',  # ZIP-based Office formats
        b'\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1': 'ole',  # OLE format
        b'\x50\x4B\x05\x06': 'zip_empty',  # Empty ZIP
        b'\x50\x4B\x07\x08': 'zip_spanned',  # Spanned ZIP
        
        # Image formats
        b'\xFF\xD8\xFF': 'jpeg',
        b'\x89\x50\x4E\x47\x0D\x0A\x1A\x0A': 'png',
        b'\x47\x49\x46\x38': 'gif',
        b'\x42\x4D': 'bmp',
        
        # PDF
        b'\x25\x50\x44\x46': 'pdf',
    }
    
    @classmethod
    def detect_format(cls, file_path: Union[str, Path]) -> Tuple[str, float]:
        """
        Detect file format with confidence score.
        Returns (format, confidence) where confidence is 0.0-1.0
        """
        file_path = Path(file_path)
        
        try:
            with open(file_path, 'rb') as f:
                header = f.read(16)
            
            # Check signatures
            for signature, format_type in cls.SIGNATURES.items():
                if header.startswith(signature):
                    confidence = 0.9
                    
                    # Additional validation for Office formats
                    if format_type == 'office_xml':
                        confidence = cls._validate_office_xml(file_path)
                    elif format_type == 'ole':
                        confidence = cls._validate_ole(file_path)
                    
                    return format_type, confidence
            
            # Fallback to extension-based detection
            ext = file_path.suffix.lower()
            mime_type, _ = mimetypes.guess_type(str(file_path))
            
            if ext in ['.pptx', '.docx', '.xlsx']:
                return 'office_xml', 0.7
            elif ext in ['.ppt', '.doc', '.xls']:
                return 'ole', 0.7
            elif ext in ['.vsd', '.vsdx']:
                return 'visio', 0.6
            elif mime_type and mime_type.startswith('image/'):
                return 'image', 0.5
            
            return 'unknown', 0.0
            
        except Exception as e:
            logging.error(f"Error detecting format for {file_path}: {e}")
            return 'unknown', 0.0
    
    @classmethod
    def _validate_office_xml(cls, file_path: Path) -> float:
        """Validate Office XML format by checking internal structure."""
        try:
            with zipfile.ZipFile(file_path, 'r') as zf:
                files = zf.namelist()
                
                # Check for required Office XML files
                required_files = ['[Content_Types].xml', '_rels/.rels']
                if all(f in files for f in required_files):
                    return 0.95
                elif any(f.startswith('ppt/') for f in files):
                    return 0.9  # PowerPoint
                elif any(f.startswith('word/') for f in files):
                    return 0.9  # Word
                elif any(f.startswith('xl/') for f in files):
                    return 0.9  # Excel
                else:
                    return 0.7
        except:
            return 0.3
    
    @classmethod
    def _validate_ole(cls, file_path: Path) -> float:
        """Validate OLE format."""
        if not OLE_AVAILABLE:
            return 0.5
        
        try:
            if olefile.isOleFile(str(file_path)):
                return 0.9
        except:
            pass
        return 0.3


class PowerPointParser:
    """Enhanced PowerPoint parser with robust error handling."""
    
    def __init__(self):
        self.logger = logging.getLogger(__name__)
    
    def parse(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """Parse PowerPoint file with multiple fallback strategies."""
        file_path = Path(file_path)
        
        # Try different parsing strategies in order of preference
        strategies = [
            self._parse_with_pptx,
            self._parse_as_zip,
            self._parse_metadata_only,
            self._parse_basic_info
        ]
        
        last_error = None
        for strategy in strategies:
            try:
                result = strategy(file_path)
                if result and result.get('success'):
                    return result
            except Exception as e:
                last_error = e
                self.logger.warning(f"Strategy {strategy.__name__} failed: {e}")
                continue
        
        # All strategies failed
        return {
            'success': False,
            'error': f"All parsing strategies failed. Last error: {last_error}",
            'file_path': str(file_path),
            'file_size': file_path.stat().st_size if file_path.exists() else 0
        }
    
    def _parse_with_pptx(self, file_path: Path) -> Dict[str, Any]:
        """Parse using python-pptx library."""
        if not PPTX_AVAILABLE:
            raise BinaryFormatError("python-pptx not available")
        
        try:
            prs = Presentation(str(file_path))
            
            slides_data = []
            total_text_length = 0
            
            for i, slide in enumerate(prs.slides):
                slide_data = {
                    'slide_number': i + 1,
                    'shapes': len(slide.shapes),
                    'text_content': [],
                    'images': 0,
                    'tables': 0
                }
                
                # Extract text from shapes
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, 'text') and shape.text:
                            slide_data['text_content'].append(shape.text.strip())
                            total_text_length += len(shape.text)
                        
                        # Count different shape types
                        if hasattr(shape, 'shape_type'):
                            if 'PICTURE' in str(shape.shape_type):
                                slide_data['images'] += 1
                            elif 'TABLE' in str(shape.shape_type):
                                slide_data['tables'] += 1
                    except Exception as e:
                        self.logger.debug(f"Error processing shape: {e}")
                        continue
                
                slides_data.append(slide_data)
            
            return {
                'success': True,
                'parser': 'python-pptx',
                'file_path': str(file_path),
                'file_size': file_path.stat().st_size,
                'slide_count': len(prs.slides),
                'total_text_length': total_text_length,
                'slides': slides_data,
                'core_properties': self._extract_core_properties(prs)
            }
            
        except (PackageNotFoundError, PresentationError) as e:
            raise BinaryFormatError(f"Invalid PowerPoint file: {e}")
        except Exception as e:
            raise BinaryFormatError(f"Error parsing with python-pptx: {e}")
    
    def _parse_as_zip(self, file_path: Path) -> Dict[str, Any]:
        """Parse PowerPoint file as ZIP archive."""
        try:
            with zipfile.ZipFile(file_path, 'r') as zf:
                files = zf.namelist()
                
                # Extract basic structure information
                ppt_files = [f for f in files if f.startswith('ppt/')]
                slide_files = [f for f in files if f.startswith('ppt/slides/slide')]
                media_files = [f for f in files if f.startswith('ppt/media/')]
                
                # Try to extract some text content
                text_content = []
                for slide_file in slide_files[:5]:  # Limit to first 5 slides
                    try:
                        with zf.open(slide_file) as sf:
                            content = sf.read().decode('utf-8', errors='ignore')
                            # Simple text extraction from XML
                            if '<a:t>' in content:
                                import re
                                texts = re.findall(r'<a:t>(.*?)</a:t>', content)
                                text_content.extend(texts)
                    except Exception as e:
                        self.logger.debug(f"Error extracting text from {slide_file}: {e}")
                
                return {
                    'success': True,
                    'parser': 'zip_extraction',
                    'file_path': str(file_path),
                    'file_size': file_path.stat().st_size,
                    'total_files': len(files),
                    'ppt_files': len(ppt_files),
                    'slide_count': len(slide_files),
                    'media_files': len(media_files),
                    'extracted_text': text_content[:10],  # First 10 text elements
                    'structure': {
                        'has_slides': len(slide_files) > 0,
                        'has_media': len(media_files) > 0,
                        'has_themes': any('theme' in f for f in files)
                    }
                }
                
        except zipfile.BadZipFile:
            raise BinaryFormatError("File is not a valid ZIP archive")
        except Exception as e:
            raise BinaryFormatError(f"Error parsing as ZIP: {e}")
    
    def _parse_metadata_only(self, file_path: Path) -> Dict[str, Any]:
        """Extract only metadata and basic file information."""
        try:
            stat = file_path.stat()
            
            # Try to read file header for additional info
            with open(file_path, 'rb') as f:
                header = f.read(1024)
            
            # Look for text patterns in header
            text_indicators = [
                b'Microsoft Office PowerPoint',
                b'PowerPoint',
                b'ppt/',
                b'slides/',
                b'presentation'
            ]
            
            confidence = sum(1 for indicator in text_indicators if indicator in header)
            
            return {
                'success': True,
                'parser': 'metadata_only',
                'file_path': str(file_path),
                'file_size': stat.st_size,
                'modified_time': stat.st_mtime,
                'confidence_score': confidence / len(text_indicators),
                'header_analysis': {
                    'likely_powerpoint': confidence > 0,
                    'header_size': len(header)
                }
            }
            
        except Exception as e:
            raise BinaryFormatError(f"Error extracting metadata: {e}")
    
    def _parse_basic_info(self, file_path: Path) -> Dict[str, Any]:
        """Last resort: basic file information only."""
        try:
            stat = file_path.stat()
            
            return {
                'success': True,
                'parser': 'basic_info',
                'file_path': str(file_path),
                'file_size': stat.st_size,
                'file_extension': file_path.suffix,
                'modified_time': stat.st_mtime,
                'warning': 'Only basic file information could be extracted'
            }
            
        except Exception as e:
            raise BinaryFormatError(f"Error getting basic info: {e}")
    
    def _extract_core_properties(self, presentation) -> Dict[str, Any]:
        """Extract core properties from presentation."""
        try:
            props = presentation.core_properties
            return {
                'title': getattr(props, 'title', None),
                'author': getattr(props, 'author', None),
                'subject': getattr(props, 'subject', None),
                'created': str(getattr(props, 'created', None)),
                'modified': str(getattr(props, 'modified', None)),
                'category': getattr(props, 'category', None),
                'comments': getattr(props, 'comments', None)
            }
        except Exception as e:
            self.logger.debug(f"Error extracting core properties: {e}")
            return {}


class VisioParser:
    """Visio file parser with fallback mechanisms."""
    
    def __init__(self):
        self.logger = logging.getLogger(__name__)
    
    def parse(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """Parse Visio file with multiple strategies."""
        file_path = Path(file_path)
        
        # Determine Visio format
        if file_path.suffix.lower() == '.vsdx':
            return self._parse_vsdx(file_path)
        elif file_path.suffix.lower() in ['.vsd', '.vss', '.vst']:
            return self._parse_legacy_visio(file_path)
        else:
            return self._parse_unknown_visio(file_path)
    
    def _parse_vsdx(self, file_path: Path) -> Dict[str, Any]:
        """Parse VSDX (XML-based) Visio files."""
        try:
            with zipfile.ZipFile(file_path, 'r') as zf:
                files = zf.namelist()
                
                # Look for Visio-specific files
                visio_files = [f for f in files if f.startswith('visio/')]
                page_files = [f for f in files if 'page' in f.lower()]
                master_files = [f for f in files if 'master' in f.lower()]
                
                # Try to extract some metadata
                metadata = {}
                if 'docProps/app.xml' in files:
                    try:
                        with zf.open('docProps/app.xml') as f:
                            content = f.read().decode('utf-8')
                            # Simple XML parsing for metadata
                            if '<Application>' in content:
                                import re
                                app_match = re.search(r'<Application>(.*?)</Application>', content)
                                if app_match:
                                    metadata['application'] = app_match.group(1)
                    except Exception as e:
                        self.logger.debug(f"Error extracting app metadata: {e}")
                
                return {
                    'success': True,
                    'parser': 'vsdx_zip',
                    'file_path': str(file_path),
                    'file_size': file_path.stat().st_size,
                    'total_files': len(files),
                    'visio_files': len(visio_files),
                    'page_files': len(page_files),
                    'master_files': len(master_files),
                    'metadata': metadata,
                    'structure': {
                        'has_pages': len(page_files) > 0,
                        'has_masters': len(master_files) > 0,
                        'has_themes': any('theme' in f for f in files)
                    }
                }
                
        except zipfile.BadZipFile:
            return {
                'success': False,
                'error': 'File is not a valid VSDX (ZIP) archive',
                'file_path': str(file_path)
            }
        except Exception as e:
            return {
                'success': False,
                'error': f'Error parsing VSDX file: {e}',
                'file_path': str(file_path)
            }
    
    def _parse_legacy_visio(self, file_path: Path) -> Dict[str, Any]:
        """Parse legacy Visio files (.vsd, .vss, .vst)."""
        try:
            if OLE_AVAILABLE and olefile.isOleFile(str(file_path)):
                with olefile.OleFileIO(str(file_path)) as ole:
                    streams = ole.listdir()
                    
                    # Look for Visio-specific streams
                    visio_streams = [s for s in streams if any(
                        keyword in str(s).lower() 
                        for keyword in ['visio', 'page', 'master', 'shape']
                    )]
                    
                    return {
                        'success': True,
                        'parser': 'ole_legacy',
                        'file_path': str(file_path),
                        'file_size': file_path.stat().st_size,
                        'total_streams': len(streams),
                        'visio_streams': len(visio_streams),
                        'streams': streams[:10],  # First 10 streams
                        'format': 'legacy_visio'
                    }
            else:
                # Fallback to basic analysis
                return self._parse_basic_binary(file_path, 'legacy_visio')
                
        except Exception as e:
            return {
                'success': False,
                'error': f'Error parsing legacy Visio file: {e}',
                'file_path': str(file_path)
            }
    
    def _parse_unknown_visio(self, file_path: Path) -> Dict[str, Any]:
        """Parse unknown Visio format."""
        return self._parse_basic_binary(file_path, 'unknown_visio')
    
    def _parse_basic_binary(self, file_path: Path, format_type: str) -> Dict[str, Any]:
        """Basic binary file analysis."""
        try:
            stat = file_path.stat()
            
            with open(file_path, 'rb') as f:
                header = f.read(1024)
            
            # Look for Visio indicators
            visio_indicators = [
                b'Visio',
                b'Microsoft Visio',
                b'vsd',
                b'Shape',
                b'Page'
            ]
            
            confidence = sum(1 for indicator in visio_indicators if indicator in header)
            
            return {
                'success': True,
                'parser': 'basic_binary',
                'file_path': str(file_path),
                'file_size': stat.st_size,
                'format_type': format_type,
                'confidence_score': confidence / len(visio_indicators),
                'header_analysis': {
                    'likely_visio': confidence > 0,
                    'header_size': len(header)
                }
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': f'Error in basic binary analysis: {e}',
                'file_path': str(file_path)
            }


class RobustBinaryParser:
    """Main binary parser with format detection and error recovery."""
    
    def __init__(self):
        self.logger = logging.getLogger(__name__)
        self.powerpoint_parser = PowerPointParser()
        self.visio_parser = VisioParser()
        self.validator = FormatValidator()
    
    def parse_file(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """Parse binary file with automatic format detection."""
        file_path = Path(file_path)
        
        if not file_path.exists():
            return {
                'success': False,
                'error': f'File not found: {file_path}',
                'file_path': str(file_path)
            }
        
        # Detect format
        format_type, confidence = self.validator.detect_format(file_path)
        
        try:
            # Route to appropriate parser
            if format_type in ['office_xml', 'ole'] or file_path.suffix.lower() in ['.ppt', '.pptx']:
                result = self.powerpoint_parser.parse(file_path)
            elif format_type == 'visio' or file_path.suffix.lower() in ['.vsd', '.vsdx', '.vss', '.vst']:
                result = self.visio_parser.parse(file_path)
            else:
                result = self._parse_generic_binary(file_path)
            
            # Add format detection info
            result['format_detection'] = {
                'detected_format': format_type,
                'confidence': confidence,
                'file_extension': file_path.suffix
            }
            
            return result
            
        except Exception as e:
            self.logger.error(f"Error parsing {file_path}: {e}")
            return {
                'success': False,
                'error': f'Parsing failed: {e}',
                'file_path': str(file_path),
                'format_detection': {
                    'detected_format': format_type,
                    'confidence': confidence
                },
                'traceback': traceback.format_exc()
            }
    
    def _parse_generic_binary(self, file_path: Path) -> Dict[str, Any]:
        """Generic binary file parser for unsupported formats."""
        try:
            stat = file_path.stat()
            
            with open(file_path, 'rb') as f:
                header = f.read(1024)
                # Sample middle and end of file
                f.seek(stat.st_size // 2)
                middle = f.read(512)
                f.seek(max(0, stat.st_size - 512))
                tail = f.read(512)
            
            # Basic analysis
            total_sample = header + middle + tail
            text_chars = sum(1 for b in total_sample if 32 <= b <= 126)
            text_ratio = text_chars / len(total_sample) if total_sample else 0
            
            return {
                'success': True,
                'parser': 'generic_binary',
                'file_path': str(file_path),
                'file_size': stat.st_size,
                'analysis': {
                    'text_ratio': text_ratio,
                    'likely_binary': text_ratio < 0.3,
                    'header_preview': header[:100].hex(),
                    'sample_size': len(total_sample)
                }
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': f'Generic parsing failed: {e}',
                'file_path': str(file_path)
            }
    
    def batch_parse(self, file_paths: List[Union[str, Path]]) -> List[Dict[str, Any]]:
        """Parse multiple files with error isolation."""
        results = []
        
        for file_path in file_paths:
            try:
                result = self.parse_file(file_path)
                results.append(result)
            except Exception as e:
                self.logger.error(f"Error in batch parsing {file_path}: {e}")
                results.append({
                    'success': False,
                    'error': f'Batch parsing error: {e}',
                    'file_path': str(file_path)
                })
        
        return results


# Example usage and testing
if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO)
    
    parser = RobustBinaryParser()
    
    # Example usage
    # result = parser.parse_file("example.pptx")
    # print(json.dumps(result, indent=2, default=str))
    
    print("Robust Binary Parser initialized successfully!")
